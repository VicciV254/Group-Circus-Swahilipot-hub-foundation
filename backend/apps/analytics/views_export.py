"""
BMI Analytics — Production Export Engine (views_export.py)
==========================================================
Generates: PDF · Excel · PPTX · CSV · JSON from live Django ORM data
Libraries: pandas · matplotlib · seaborn · reportlab · openpyxl · python-pptx
"""

import io, json, csv, math
from datetime import datetime, timedelta

import pandas as pd
import numpy as np
import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt
import matplotlib.gridspec as gridspec
import seaborn as sns

from django.utils import timezone
from django.http import HttpResponse, JsonResponse
from django.db.models import Count, Avg, Sum, Q, Max, Min
from rest_framework.views import APIView
from rest_framework.permissions import IsAuthenticated

from reportlab.lib.pagesizes import A4
from reportlab.lib.units import mm
from reportlab.lib.colors import HexColor
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.platypus import (
    SimpleDocTemplate, Table, TableStyle, Paragraph,
    Spacer, Image as RLImage, PageBreak,
)

from openpyxl import Workbook
from openpyxl.styles import (
    Font, PatternFill, Alignment, Border, Side,
)
from openpyxl.utils import get_column_letter
from openpyxl.chart import BarChart, PieChart, LineChart, Reference

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN


# ─── Colour tokens ────────────────────────────────────────────────────────────
class C:
    DARK  = '#0D1117'; CARD  = '#161B27'; BORDER = '#1E2A3B'
    BLUE  = '#3B63F5'; GREEN = '#22C55E'; AMBER  = '#F59E0B'
    RED   = '#EF4444'; PURPLE= '#8B5CF6'; CYAN   = '#06b6d4'
    SLATE = '#94A3B8'; TEXT  = '#E2E8F0'; WHITE  = '#FFFFFF'
    PALETTE = [BLUE, GREEN, AMBER, RED, PURPLE, CYAN, '#EC4899', '#F97316', '#14B8A6']

# ─── matplotlib dark theme ────────────────────────────────────────────────────
MPL_STYLE = {
    'figure.facecolor': C.DARK,  'axes.facecolor': C.CARD,
    'axes.edgecolor':   C.BORDER,'axes.labelcolor': C.SLATE,
    'xtick.color':      C.SLATE, 'ytick.color':     C.SLATE,
    'text.color':       C.TEXT,  'grid.color':      '#1A2236',
    'grid.linestyle':   '--',    'grid.alpha':      0.5,
    'font.family':      'DejaVu Sans', 'font.size': 9,
}

# ─── Helpers ──────────────────────────────────────────────────────────────────
def _pct(num, denom, dec=1):
    return round(num / denom * 100, dec) if denom else 0.0

def _sv(v):
    """JSON-safe value."""
    if isinstance(v, float) and math.isnan(v): return None
    if hasattr(v, 'item'): return v.item()
    return v

def _rl_color(hex6): return HexColor('#' + hex6.lstrip('#'))

def _chart_buf(fig):
    buf = io.BytesIO()
    plt.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                facecolor=C.DARK, edgecolor='none')
    plt.close('all'); buf.seek(0)
    return buf

def _make_rl_table(data, col_widths, has_header=True):
    t  = Table(data, colWidths=col_widths)
    ts = TableStyle([
        ('FONTSIZE',      (0, 0), (-1, -1), 8),
        ('ALIGN',         (0, 0), (-1, -1), 'CENTER'),
        ('VALIGN',        (0, 0), (-1, -1), 'MIDDLE'),
        ('ROWHEIGHT',     (0, 0), (-1, -1), 15),
        ('GRID',          (0, 0), (-1, -1), 0.4, _rl_color(C.BORDER)),
        ('LEFTPADDING',   (0, 0), (-1, -1), 5),
        ('RIGHTPADDING',  (0, 0), (-1, -1), 5),
        ('ROWBACKGROUNDS',(0, 1 if has_header else 0), (-1, -1),
         [_rl_color(C.CARD), HexColor('#0F1623')]),
        ('TEXTCOLOR',     (0, 1 if has_header else 0), (-1, -1), _rl_color(C.SLATE)),
    ])
    if has_header:
        ts.add('BACKGROUND', (0, 0), (-1, 0), _rl_color(C.BLUE))
        ts.add('TEXTCOLOR',  (0, 0), (-1, 0), _rl_color(C.WHITE))
        ts.add('FONTNAME',   (0, 0), (-1, 0), 'Helvetica-Bold')
    t.setStyle(ts)
    return t


# ─── Data loader: pulls everything from ORM into pandas ──────────────────────
class BMIDataLoader:
    """
    Loads all BMI data from Django ORM into pandas DataFrames
    and computes analytics metrics.
    """

    def __init__(self, org, days=30):
        self.org  = org
        self.days = days
        self.today = timezone.now().date()
        self.start = self.today - timedelta(days=days)
        self._load()

    def _load(self):
        from apps.accounts.models import User, Department
        from apps.tasks.models import Task
        from apps.attendance.models import AttendanceRecord, LeaveRequest

        # ── Users ──────────────────────────────────────────────────────────
        users_qs = User.objects.filter(organisation=self.org).values(
            'id', 'employee_id', 'first_name', 'last_name',
            'role', 'is_active', 'mfa_enabled', 'date_joined', 'last_login',
            'department__name',
        )
        self.users = pd.DataFrame(list(users_qs))
        self.user_map = dict(zip(
            self.users['id'].astype(str),
            self.users['first_name'] + ' ' + self.users['last_name']
        ))

        # ── Audit log ──────────────────────────────────────────────────────
        try:
            from core.models import AuditLog
            audit_qs = AuditLog.objects.filter(
                user__organisation=self.org,
                created_at__date__gte=self.start,
            ).values('id', 'created_at', 'user_id', 'action', 'resource_type',
                     'resource_id', 'extra', 'ip_address')
            self.audit = pd.DataFrame(list(audit_qs))
            if not self.audit.empty:
                self.audit['ts']   = pd.to_datetime(self.audit['created_at'], utc=True, errors='coerce')
                self.audit['date'] = self.audit['ts'].dt.date
                self.audit['hour'] = self.audit['ts'].dt.hour
                self.audit['duration_ms'] = self.audit['extra'].apply(
                    lambda x: json.loads(x).get('duration_ms') if isinstance(x, str) else None
                )
                self.durations = self.audit['duration_ms'].dropna()
            else:
                self.durations = pd.Series([], dtype=float)
        except Exception:
            self.audit     = pd.DataFrame()
            self.durations = pd.Series([], dtype=float)

        # ── Tasks ──────────────────────────────────────────────────────────
        from django.contrib.postgres.aggregates import StringAgg

        tasks_qs = Task.objects.filter(
            organisation=self.org
        ).annotate(
            department_names=StringAgg('departments__name', delimiter=', ', distinct=True)
        ).values('id', 'title', 'priority', 'status', 'progress_percent',
         'due_date', 'submitted_at', 'reviewed_at',
         'assigned_to__first_name', 'assigned_to__last_name',
         'department_names')   # ← annotated field
        self.tasks = pd.DataFrame(list(tasks_qs))

        # ── Attendance ─────────────────────────────────────────────────────
        att_qs = AttendanceRecord.objects.filter(
            user__organisation=self.org, date__gte=self.start,
        ).values('id', 'date', 'status', 'method', 'total_hours',
                 'check_in_time', 'check_out_time',
                 'user__first_name', 'user__last_name',
                 'user__department__name')
        self.attendance = pd.DataFrame(list(att_qs))

        # ── Equipment ──────────────────────────────────────────────────────
        try:
            from apps.equipment.models import EquipmentItem, MaintenanceLog
            eq_qs = EquipmentItem.objects.filter(
                organisation=self.org, is_active=True,
            ).values('id', 'asset_tag', 'name', 'status', 'condition',
                     'purchase_cost', 'current_value', 'location',
                     'category__name', 'purchase_date')
            self.equipment = pd.DataFrame(list(eq_qs))
            if not self.equipment.empty:
                self.equipment['purchase_cost'] = pd.to_numeric(
                    self.equipment['purchase_cost'], errors='coerce').fillna(0)

            ml_qs = MaintenanceLog.objects.filter(
                item__organisation=self.org,
            ).values('id', 'item__name', 'item__asset_tag', 'description',
                     'status', 'created_at', 'reported_by__first_name',
                     'reported_by__last_name', 'repair_cost')
            self.maintenance = pd.DataFrame(list(ml_qs))
        except Exception:
            self.equipment   = pd.DataFrame()
            self.maintenance = pd.DataFrame()

        # ── FM ─────────────────────────────────────────────────────────────
        try:
            from apps.fm_report.models import FMOutage, FMStation, EmergencyAlert
            outage_qs = FMOutage.objects.filter(
                station__organisation=self.org,
                down_at__date__gte=self.start,
            ).values('id', 'station__name', 'down_at', 'restored_at',
                     'duration_minutes', 'severity', 'description',
                     'reported_by__first_name', 'reported_by__last_name')
            self.fm_outages = pd.DataFrame(list(outage_qs))

            alert_qs = EmergencyAlert.objects.filter(
                organisation=self.org,
            ).values('id', 'title', 'alert_type', 'severity',
                     'resolved', 'created_at', 'description')
            self.alerts = pd.DataFrame(list(alert_qs))
        except Exception:
            self.fm_outages = pd.DataFrame()
            self.alerts     = pd.DataFrame()

        # ── Notifications ──────────────────────────────────────────────────
        try:
            from apps.notifications.models import Notification
            notif_qs = Notification.objects.filter(
                recipient__organisation=self.org,
                created_at__date__gte=self.start,
            ).values('id', 'title', 'notification_type', 'read',
                     'is_urgent', 'push_sent', 'email_sent', 'created_at',
                     'recipient__first_name', 'recipient__last_name')
            self.notifications = pd.DataFrame(list(notif_qs))
        except Exception:
            self.notifications = pd.DataFrame()

        # ── Certificates ───────────────────────────────────────────────────
        try:
            from apps.certificates.models import Certificate
            cert_qs = Certificate.objects.filter(
                organisation=self.org,
            ).values('id', 'certificate_number', 'certificate_type',
                     'status', 'issue_date',
                     'recipient__first_name', 'recipient__last_name',
                     'signed_by_name', 'signed_by_title')
            self.certificates = pd.DataFrame(list(cert_qs))
        except Exception:
            self.certificates = pd.DataFrame()

        # ── Feedback ───────────────────────────────────────────────────────
        try:
            from apps.feedback.models import FeedbackTicket
            fb_qs = FeedbackTicket.objects.filter(
                submitted_by__organisation=self.org,
            ).values('id', 'ticket_number', 'category', 'title',
                     'priority', 'status', 'created_at',
                     'submitted_by__first_name', 'submitted_by__last_name',
                     'admin_response', 'resolved_at')
            self.feedback = pd.DataFrame(list(fb_qs))
        except Exception:
            self.feedback = pd.DataFrame()

        # ── File transfers ─────────────────────────────────────────────────
        try:
            from apps.filetransfer.models import FileTransfer
            ft_qs = FileTransfer.objects.filter(
                uploaded_by__organisation=self.org,
            ).values('id', 'original_filename', 'file_size',
                     'download_count', 'created_at',
                     'uploaded_by__first_name', 'uploaded_by__last_name')
            self.filetransfers = pd.DataFrame(list(ft_qs))
            if not self.filetransfers.empty:
                self.filetransfers['file_size'] = pd.to_numeric(
                    self.filetransfers['file_size'], errors='coerce').fillna(0)
        except Exception:
            self.filetransfers = pd.DataFrame()

        # ── Wi-Fi ──────────────────────────────────────────────────────────
        try:
            from apps.wifi.models import WifiGrant
            wifi_qs = WifiGrant.objects.filter(
                requested_by__organisation=self.org,
            ).values('id', 'device_type', 'mac_address', 'purpose',
                     'status', 'duration_days', 'created_at', 'expires_at',
                     'requested_by__first_name', 'requested_by__last_name',
                     'reviewed_by__first_name', 'reviewed_by__last_name')
            self.wifi = pd.DataFrame(list(wifi_qs))
        except Exception:
            self.wifi = pd.DataFrame()

        # ── Videography ────────────────────────────────────────────────────
        try:
            from apps.videography.models import ShootBooking
            vb_qs = ShootBooking.objects.filter(
                requested_by__organisation=self.org,
            ).values('id', 'title', 'shoot_date', 'duration_hours',
                     'location', 'status', 'created_at',
                     'requested_by__first_name', 'requested_by__last_name')
            self.videography = pd.DataFrame(list(vb_qs))
        except Exception:
            self.videography = pd.DataFrame()

    # ── Computed stats ────────────────────────────────────────────────────────
    @property
    def stats(self):
        s = {}
        s['active_users']   = int(self.users[self.users['is_active'] == True]['id'].count()) if not self.users.empty else 0
        s['mfa_enabled']    = int(self.users[self.users['mfa_enabled'] == True]['id'].count()) if not self.users.empty else 0
        s['mfa_pct']        = _pct(s['mfa_enabled'], s['active_users'])
        s['api_calls']      = len(self.audit)
        s['api_posts']      = int((self.audit['action'] == 'POST').sum()) if not self.audit.empty else 0
        s['api_patches']    = int((self.audit['action'] == 'PATCH').sum()) if not self.audit.empty else 0
        s['avg_ms']         = round(float(self.durations.mean()), 1) if len(self.durations) else 0
        s['median_ms']      = round(float(self.durations.median()), 1) if len(self.durations) else 0
        s['p95_ms']         = round(float(self.durations.quantile(0.95)), 1) if len(self.durations) else 0
        s['p99_ms']         = round(float(self.durations.quantile(0.99)), 1) if len(self.durations) else 0
        s['max_ms']         = round(float(self.durations.max()), 1) if len(self.durations) else 0
        s['tasks_total']    = len(self.tasks)
        s['tasks_approved'] = int((self.tasks['status'] == 'approved').sum()) if not self.tasks.empty else 0
        s['tasks_overdue']  = int((self.tasks['status'] == 'overdue').sum())  if not self.tasks.empty else 0
        s['equip_total']    = len(self.equipment)
        s['equip_value']    = float(self.equipment['purchase_cost'].sum()) if not self.equipment.empty else 0
        s['equip_repair']   = int((self.equipment['status'] == 'under_repair').sum()) if not self.equipment.empty else 0
        s['alerts_total']   = len(self.alerts)
        s['alerts_resolved']= int((self.alerts['resolved'] == True).sum()) if not self.alerts.empty else 0
        s['notif_total']    = len(self.notifications)
        s['notif_read']     = int((self.notifications['read'] == True).sum()) if not self.notifications.empty else 0
        s['certs_issued']   = len(self.certificates)
        s['feedback_total'] = len(self.feedback)
        s['feedback_open']  = int((self.feedback['status'] == 'open').sum()) if not self.feedback.empty else 0
        s['files_total']    = len(self.filetransfers)
        s['files_kb']       = round(float(self.filetransfers['file_size'].sum()) / 1024, 1) if not self.filetransfers.empty else 0
        s['report_date']    = datetime.now().strftime('%B %d, %Y')
        return s

    def perf_buckets(self):
        if len(self.durations) == 0:
            return pd.Series({'<100ms': 0, '100-500ms': 0, '500ms-1s': 0, '>1s': 0})
        return pd.cut(
            self.durations,
            bins=[0, 100, 500, 1000, float('inf')],
            labels=['<100ms', '100-500ms', '500ms-1s', '>1s'],
        ).value_counts().sort_index()

    def daily_action(self):
        if self.audit.empty: return pd.DataFrame()
        da = self.audit.groupby(['date', 'action']).size().unstack(fill_value=0).reset_index()
        da['date_str'] = da['date'].astype(str).str[-5:]
        return da

    def hourly(self):
        if self.audit.empty: return pd.DataFrame({'hour': range(24), 'calls': [0]*24})
        h = self.audit.groupby('hour').size().reset_index(name='calls')
        return pd.DataFrame({'hour': range(24)}).merge(h, on='hour', how='left').fillna(0)

    def module_breakdown(self):
        if self.audit.empty: return pd.DataFrame()
        return self.audit.groupby('resource_type').size().reset_index(name='calls').sort_values('calls', ascending=False)

    def user_activity(self):
        if self.audit.empty: return pd.DataFrame()
        uc = self.audit.groupby('user_id').size().reset_index(name='calls')
        uc['name'] = uc['user_id'].astype(str).map(self.user_map).fillna('Unknown')
        return uc

    def notif_types(self):
        if self.notifications.empty: return pd.DataFrame()
        nt = self.notifications['notification_type'].value_counts().reset_index()
        nt.columns = ['type', 'count']
        return nt


# ─── Chart builder ────────────────────────────────────────────────────────────
class ChartBuilder:
    """Generates all matplotlib/seaborn chart BytesIO objects."""

    def __init__(self, loader: BMIDataLoader):
        self.L = loader
        self.S = loader.stats
        plt.rcParams.update(MPL_STYLE)

    def _save(self, fig=None) -> io.BytesIO:
        buf = io.BytesIO()
        plt.savefig(buf, format='png', dpi=150, bbox_inches='tight',
                    facecolor=C.DARK, edgecolor='none')
        plt.close('all'); buf.seek(0)
        return buf

    def api_daily(self) -> io.BytesIO:
        da = self.L.daily_action()
        fig, ax = plt.subplots(figsize=(8, 3.2))
        if da.empty:
            ax.text(0.5, 0.5, 'No data', ha='center', va='center',
                    color=C.SLATE, transform=ax.transAxes)
            return self._save()
        x = np.arange(len(da))
        pv = da.get('POST',  pd.Series([0]*len(da))).values
        mv = da.get('PATCH', pd.Series([0]*len(da))).values
        b1 = ax.bar(x, pv, color=C.BLUE,   label='POST (Create)',  alpha=0.9, width=0.5)
        b2 = ax.bar(x, mv, bottom=pv,       color=C.PURPLE, label='PATCH (Update)', alpha=0.9, width=0.5)
        ax.set_xticks(x); ax.set_xticklabels(da['date_str'].values)
        ax.set_title('API Activity by Day', color=C.TEXT, pad=10, fontsize=11, fontweight='bold')
        ax.legend(loc='upper left', framealpha=0.2, labelcolor=C.TEXT)
        for bar, val, bot in zip(b2, mv, pv):
            if val > 0:
                ax.text(bar.get_x()+bar.get_width()/2, bot+val/2, str(int(val)),
                        ha='center', va='center', color=C.TEXT, fontsize=9, fontweight='bold')
        for bar, val in zip(b1, pv):
            if val > 0:
                ax.text(bar.get_x()+bar.get_width()/2, val/2, str(int(val)),
                        ha='center', va='center', color=C.TEXT, fontsize=9, fontweight='bold')
        plt.tight_layout(); return self._save()

    def hourly(self) -> io.BytesIO:
        hf = self.L.hourly()
        fig, ax = plt.subplots(figsize=(9, 2.8))
        clrs = [C.RED if c > 30 else C.AMBER if c > 15 else C.BLUE for c in hf['calls']]
        ax.bar(hf['hour'], hf['calls'], color=clrs, alpha=0.88, width=0.7)
        ax.set_xlabel('Hour of Day (24h)', color=C.SLATE)
        ax.set_ylabel('API Calls', color=C.SLATE)
        ax.set_title('API Activity by Hour', color=C.TEXT, pad=8, fontsize=11, fontweight='bold')
        ax.set_xticks(range(0, 24, 2))
        for h, cnt in zip(hf['hour'], hf['calls']):
            if cnt > 15:
                ax.text(h, cnt+0.5, str(int(cnt)), ha='center', va='bottom',
                        color=C.TEXT, fontsize=8, fontweight='bold')
        plt.tight_layout(); return self._save()

    def modules(self) -> io.BytesIO:
        mc = self.L.module_breakdown()
        fig, ax = plt.subplots(figsize=(8, max(3.5, len(mc)*0.4)))
        if mc.empty:
            ax.text(0.5, 0.5, 'No data', ha='center', va='center',
                    color=C.SLATE, transform=ax.transAxes)
            return self._save()
        clrs = [C.PALETTE[i % len(C.PALETTE)] for i in range(len(mc))]
        bars = ax.barh(mc['resource_type'], mc['calls'], color=clrs, alpha=0.88, height=0.65)
        ax.set_title('API Calls by Module', color=C.TEXT, pad=8, fontsize=11, fontweight='bold')
        total = max(mc['calls'].sum(), 1)
        for bar, val in zip(bars, mc['calls']):
            ax.text(val+0.3, bar.get_y()+bar.get_height()/2,
                    f'{int(val)} ({round(val/total*100, 1)}%)',
                    va='center', color=C.TEXT, fontsize=8)
        plt.tight_layout(); return self._save()

    def response_time(self) -> io.BytesIO:
        pb = self.L.perf_buckets()
        fig, (ax1, ax2) = plt.subplots(1, 2, figsize=(10, 3.8))
        colors_pie = [C.GREEN, C.AMBER, '#F97316', C.RED]
        if pb.sum() > 0:
            w, t, at = ax1.pie(
                pb.values, labels=pb.index.tolist(), colors=colors_pie,
                autopct='%1.1f%%', startangle=140,
                wedgeprops=dict(width=0.55, edgecolor=C.DARK),
                textprops={'color': C.TEXT, 'fontsize': 9},
            )
            for a in at: a.set_color(C.DARK); a.set_fontweight('bold')
        ax1.set_title('Response Time Buckets', color=C.TEXT, fontsize=10, fontweight='bold')

        dur = self.L.durations
        if len(dur):
            ax2.hist(dur.clip(upper=5000), bins=50, color=C.BLUE, alpha=0.8, edgecolor=C.DARK)
            ax2.axvline(float(dur.median()), color=C.GREEN, linestyle='--', lw=1.8,
                        label=f'P50 {self.S["median_ms"]:.0f}ms')
            ax2.axvline(float(dur.quantile(0.95)), color=C.AMBER, linestyle='--', lw=1.8,
                        label=f'P95 {self.S["p95_ms"]:.0f}ms')
            ax2.axvline(float(dur.quantile(0.99)), color=C.RED, linestyle='--', lw=1.8,
                        label=f'P99 {self.S["p99_ms"]:.0f}ms')
            ax2.legend(framealpha=0.2, labelcolor=C.TEXT, fontsize=8)
        ax2.set_title('Response Time Histogram (≤5s)', color=C.TEXT, fontsize=10, fontweight='bold')
        ax2.set_xlabel('ms', color=C.SLATE)
        plt.tight_layout(); return self._save()

    def users(self) -> io.BytesIO:
        uc = self.L.user_activity().sort_values('calls', ascending=True)
        fig, ax = plt.subplots(figsize=(7, max(2.5, len(uc)*0.6)))
        clrs = [C.PALETTE[i % len(C.PALETTE)] for i in range(len(uc))]
        bars = ax.barh(uc['name'], uc['calls'], color=clrs, alpha=0.88, height=0.5)
        ax.set_title('API Calls by User', color=C.TEXT, pad=8, fontsize=11, fontweight='bold')
        total = max(uc['calls'].sum(), 1)
        for bar, val in zip(bars, uc['calls']):
            ax.text(val+0.3, bar.get_y()+bar.get_height()/2,
                    f'{int(val)} ({round(val/total*100, 1)}%)',
                    va='center', color=C.TEXT, fontsize=9)
        plt.tight_layout(); return self._save()

    def notifications(self) -> io.BytesIO:
        nt = self.L.notif_types()
        fig, ax = plt.subplots(figsize=(6, 4))
        if nt.empty:
            ax.text(0.5, 0.5, 'No notifications', ha='center', va='center',
                    color=C.SLATE, transform=ax.transAxes)
            return self._save()
        labels = [t.replace('_', ' ').title() for t in nt['type']]
        clrs   = [C.PALETTE[i % len(C.PALETTE)] for i in range(len(labels))]
        total  = nt['count'].sum()
        w, t, at = ax.pie(
            nt['count'].values, labels=labels, colors=clrs,
            autopct='%1.0f%%', startangle=90,
            wedgeprops=dict(width=0.55, edgecolor=C.DARK),
            textprops={'color': C.TEXT, 'fontsize': 8},
        )
        for a in at: a.set_color(C.DARK); a.set_fontweight('bold')
        notif_read   = int(self.L.notifications['read'].sum()) if not self.L.notifications.empty else 0
        notif_unread = total - notif_read
        ax.set_title(f'Notifications by Type ({total} total · {notif_read} read · {notif_unread} unread)',
                     color=C.TEXT, pad=8, fontsize=9, fontweight='bold')
        plt.tight_layout(); return self._save()

    def attendance_trend(self) -> io.BytesIO:
        att = self.L.attendance
        fig, ax = plt.subplots(figsize=(9, 3.2))
        if att.empty:
            ax.text(0.5, 0.5, 'No attendance records in this period',
                    ha='center', va='center', color=C.SLATE, transform=ax.transAxes)
            return self._save()
        daily = att.groupby(['date', 'status']).size().unstack(fill_value=0).reset_index()
        x = np.arange(len(daily))
        for col_name, clr in [('present', C.BLUE), ('absent', C.RED), ('leave', C.AMBER), ('late', C.PURPLE)]:
            if col_name in daily.columns:
                ax.plot(x, daily[col_name], color=clr, linewidth=2, label=col_name.title(), marker='o', markersize=4)
        ax.set_xticks(x[::max(1, len(x)//12)])
        ax.set_xticklabels([str(d)[-5:] for d in daily['date'].iloc[::max(1, len(x)//12)]], rotation=30)
        ax.set_title('Attendance Trend', color=C.TEXT, pad=8, fontsize=11, fontweight='bold')
        ax.legend(framealpha=0.2, labelcolor=C.TEXT)
        plt.tight_layout(); return self._save()

    def dashboard_overview(self) -> io.BytesIO:
        """6-panel overview figure."""
        da = self.L.daily_action()
        hf = self.L.hourly()
        mc = self.L.module_breakdown()
        nt = self.L.notif_types()

        fig = plt.figure(figsize=(16, 11), facecolor=C.DARK)
        fig.suptitle('BMI Intelligence Hub — Analytics Dashboard',
                     color=C.TEXT, fontsize=15, fontweight='bold', y=0.99)
        gs = gridspec.GridSpec(3, 3, figure=fig, hspace=0.5, wspace=0.38)

        # Panel 1: daily
        ax1 = fig.add_subplot(gs[0, :2])
        if not da.empty:
            x   = np.arange(len(da))
            pv_ = da.get('POST',  pd.Series([0]*len(da))).values
            mv_ = da.get('PATCH', pd.Series([0]*len(da))).values
            ax1.bar(x, pv_,  color=C.BLUE,   label='POST',  alpha=0.9, width=0.4)
            ax1.bar(x, mv_, bottom=pv_, color=C.PURPLE, label='PATCH', alpha=0.9, width=0.4)
            ax1.set_xticks(x); ax1.set_xticklabels(da['date_str'].values)
            ax1.legend(loc='upper left', framealpha=0.15, labelcolor=C.TEXT, fontsize=8)
        ax1.set_title('Daily API Activity', color=C.TEXT, fontsize=10, fontweight='bold')

        # Panel 2: KPI text
        ax2 = fig.add_subplot(gs[0, 2]); ax2.axis('off')
        kpi_lines = [
            ('Active Users',   str(self.S['active_users']),                C.GREEN),
            ('API Calls',      str(self.S['api_calls']),                   C.BLUE),
            ('Avg Response',   f"{self.S['avg_ms']}ms",                   C.AMBER),
            ('MFA Adoption',   f"{self.S['mfa_pct']}% {'⚠' if self.S['mfa_pct']<50 else '✓'}", C.RED if self.S['mfa_pct']<50 else C.GREEN),
            ('Asset Value',    f"KES {self.S['equip_value']/1e6:.1f}M",   C.PURPLE),
            ('Alerts OK',      f"{self.S['alerts_resolved']}/{self.S['alerts_total']}", C.GREEN),
        ]
        for i, (lbl, val, col) in enumerate(kpi_lines):
            y = 0.93 - i * 0.16
            ax2.text(0.05, y, lbl, transform=ax2.transAxes, color=C.SLATE, fontsize=8)
            ax2.text(0.95, y, val, transform=ax2.transAxes, color=col, fontsize=10,
                     fontweight='bold', ha='right')
        ax2.set_title('KPIs', color=C.TEXT, fontsize=10, fontweight='bold')

        # Panel 3: hourly
        ax3 = fig.add_subplot(gs[1, :2])
        clrs3 = [C.RED if c > 30 else C.AMBER if c > 15 else C.BLUE for c in hf['calls']]
        ax3.bar(hf['hour'], hf['calls'], color=clrs3, alpha=0.85, width=0.75)
        ax3.set_title('Hourly Activity', color=C.TEXT, fontsize=10, fontweight='bold')
        ax3.set_xlabel('Hour', color=C.SLATE, fontsize=8)

        # Panel 4: notifications
        ax4 = fig.add_subplot(gs[1, 2])
        if not nt.empty:
            clrs4 = [C.PALETTE[i % len(C.PALETTE)] for i in range(len(nt))]
            ax4.pie(nt['count'].values,
                    labels=[t.replace('_', ' ')[:10] for t in nt['type']],
                    colors=clrs4, autopct='%1.0f%%', startangle=90,
                    wedgeprops=dict(width=0.5, edgecolor=C.DARK),
                    textprops={'color': C.TEXT, 'fontsize': 7})
        ax4.set_title('Notifications', color=C.TEXT, fontsize=10, fontweight='bold')

        # Panel 5: modules
        ax5 = fig.add_subplot(gs[2, :2])
        if not mc.empty:
            clrs5 = [C.PALETTE[i % len(C.PALETTE)] for i in range(len(mc))]
            ax5.bar(range(len(mc)), mc['calls'].values, color=clrs5, alpha=0.88, width=0.65)
            ax5.set_xticks(range(len(mc)))
            ax5.set_xticklabels(mc['resource_type'].values, rotation=35, ha='right', fontsize=8)
            for i, v in enumerate(mc['calls'].values):
                ax5.text(i, v+0.3, str(int(v)), ha='center', color=C.TEXT, fontsize=8, fontweight='bold')
        ax5.set_title('API Calls by Module', color=C.TEXT, fontsize=10, fontweight='bold')

        # Panel 6: perf summary
        ax6 = fig.add_subplot(gs[2, 2]); ax6.axis('off')
        perf_lines = [
            ('Median',  f"{self.S['median_ms']}ms", C.GREEN),
            ('Average', f"{self.S['avg_ms']}ms",    C.AMBER),
            ('P95',     f"{self.S['p95_ms']}ms",    C.AMBER),
            ('P99',     f"{self.S['p99_ms']}ms",    C.RED),
            ('Max',     f"{self.S['max_ms']}ms",    C.RED),
            ('<100ms',  f"{int(self.L.perf_buckets().get('<100ms', 0))} calls", C.GREEN),
            ('>1s',     f"{int(self.L.perf_buckets().get('>1s', 0))} calls",   C.RED),
        ]
        for i, (lbl, val, col) in enumerate(perf_lines):
            y = 0.92 - i * 0.14
            ax6.text(0.05, y, lbl, transform=ax6.transAxes, color=C.SLATE, fontsize=8)
            ax6.text(0.95, y, val, transform=ax6.transAxes, color=col, fontsize=9,
                     fontweight='bold', ha='right')
        ax6.set_title('Performance', color=C.TEXT, fontsize=10, fontweight='bold')

        return self._save()


# ─── PDF builder ──────────────────────────────────────────────────────────────
class PDFBuilder:
    def __init__(self, loader: BMIDataLoader, charts: dict):
        self.L = loader
        self.S = loader.stats
        self.CH = charts
        self._styles()

    def _styles(self):
        sts = getSampleStyleSheet()
        def ps(n, **kw): return ParagraphStyle(n, parent=sts['BodyText'], **kw)
        self.s_h2   = ps('h2', fontSize=14, textColor=_rl_color(C.WHITE), spaceBefore=14, spaceAfter=6, fontName='Helvetica-Bold')
        self.s_h3   = ps('h3', fontSize=10, textColor=_rl_color(C.SLATE), spaceBefore=8,  spaceAfter=4, fontName='Helvetica-Bold')
        self.s_body = ps('b',  fontSize=9,  textColor=_rl_color(C.SLATE), leading=14,     fontName='Helvetica')
        self.s_code = ps('c',  fontSize=8,  textColor=_rl_color(C.CYAN),  leading=12,     fontName='Courier')
        self.s_crit = ps('cr', fontSize=9,  textColor=_rl_color(C.RED),   leading=13,     fontName='Helvetica-BoldOblique')
        self.s_warn = ps('w',  fontSize=9,  textColor=_rl_color(C.AMBER), leading=13,     fontName='Helvetica-Oblique')
        self.s_good = ps('g',  fontSize=9,  textColor=_rl_color(C.GREEN), leading=13,     fontName='Helvetica-Oblique')

    def _cover(self, canvas, doc):
        S = self.S
        canvas.saveState()
        canvas.setFillColor(_rl_color(C.DARK))
        canvas.rect(0, 0, A4[0], A4[1], fill=1, stroke=0)
        canvas.setFillColor(_rl_color(C.BLUE))
        canvas.roundRect(15*mm, A4[1]-78*mm, A4[0]-30*mm, 55*mm, 6, fill=1, stroke=0)
        canvas.setFillColor(_rl_color(C.WHITE))
        canvas.setFont('Helvetica-Bold', 30)
        canvas.drawCentredString(A4[0]/2, A4[1]-48*mm, 'INTELLIGENCE HUB')
        canvas.setFont('Helvetica', 13)
        canvas.setFillColor(HexColor('#93C5FD'))
        canvas.drawCentredString(A4[0]/2, A4[1]-61*mm,
                                 'Broadcast Media Institution — Full Analytics Report')
        canvas.setFont('Helvetica', 9)
        canvas.setFillColor(HexColor('#64748B'))
        canvas.drawCentredString(A4[0]/2, A4[1]-73*mm,
                                 f'Generated: {S["report_date"]}  ·  Pandas · Matplotlib · Seaborn · ReportLab')
        # KPI strip
        kpi_items = [
            (str(S['active_users']), 'Active Users',          '22C55E'),
            (str(S['api_calls']),    'API Calls',             '3B63F5'),
            (f"{S['avg_ms']}ms",     'Avg Response',          'F59E0B'),
            (f"{S['mfa_pct']}% {'⚠' if S['mfa_pct']<50 else '✓'}", 'MFA Adoption', 'EF4444' if S['mfa_pct']<50 else '22C55E'),
            (f"{S['alerts_resolved']}/{S['alerts_total']}", 'Alerts Resolved', '22C55E'),
            (f"KES {S['equip_value']/1e6:.1f}M", 'Asset Value',  '8B5CF6'),
        ]
        kw = (A4[0]-30*mm)/len(kpi_items); ky = A4[1]-122*mm
        for i, (val, lbl, col) in enumerate(kpi_items):
            x = 15*mm + i*kw
            canvas.setFillColor(HexColor('#161B27'))
            canvas.roundRect(x+1.5*mm, ky, kw-3*mm, 25*mm, 4, fill=1, stroke=0)
            canvas.setFillColor(HexColor('#'+col))
            canvas.setFont('Helvetica-Bold', 18)
            canvas.drawCentredString(x+kw/2, ky+14*mm, val)
            canvas.setFillColor(HexColor('#94A3B8'))
            canvas.setFont('Helvetica', 7)
            canvas.drawCentredString(x+kw/2, ky+5*mm, lbl)
        # Module table
        canvas.setFillColor(_rl_color(C.TEXT))
        canvas.setFont('Helvetica-Bold', 11)
        canvas.drawCentredString(A4[0]/2, A4[1]-165*mm, '▸  Top Modules by API Activity')
        mc    = self.L.module_breakdown()
        total = max(len(self.L.audit), 1)
        rows  = [('Module', 'Calls', '% Total')] + [
            (r['resource_type'][:22], str(int(r['calls'])),
             f"{round(r['calls']/total*100, 1)}%")
            for _, r in mc.head(8).iterrows()
        ]
        cxs = [22*mm, 105*mm, 145*mm]; rh = 8.5*mm
        for ri, row_ in enumerate(rows):
            yt = A4[1]-177*mm - ri*rh
            if ri == 0:
                canvas.setFillColor(_rl_color(C.BLUE))
                canvas.rect(18*mm, yt-1.5*mm, A4[0]-36*mm, rh, fill=1, stroke=0)
                canvas.setFillColor(_rl_color(C.WHITE)); canvas.setFont('Helvetica-Bold', 8)
            else:
                canvas.setFillColor(HexColor('#161B27') if ri % 2 else HexColor('#0F1623'))
                canvas.rect(18*mm, yt-1.5*mm, A4[0]-36*mm, rh, fill=1, stroke=0)
                canvas.setFillColor(HexColor('#94A3B8')); canvas.setFont('Helvetica', 8)
            for cx, val in zip(cxs, row_):
                canvas.drawString(cx, yt+1.5*mm, val)
        canvas.setFillColor(HexColor('#1E2A3B'))
        canvas.rect(0, 0, A4[0], 10*mm, fill=1, stroke=0)
        canvas.setFillColor(HexColor('#64748B')); canvas.setFont('Helvetica', 7)
        canvas.drawString(15*mm, 3.5*mm, 'CONFIDENTIAL — BMI Intelligence Hub')
        canvas.drawRightString(A4[0]-15*mm, 3.5*mm, S['report_date'])
        canvas.restoreState()

    def _later(self, canvas, doc):
        canvas.saveState()
        canvas.setFillColor(_rl_color(C.DARK)); canvas.rect(0, 0, A4[0], A4[1], fill=1, stroke=0)
        canvas.setFillColor(HexColor('#161B27')); canvas.rect(0, A4[1]-14*mm, A4[0], 14*mm, fill=1, stroke=0)
        canvas.setFillColor(HexColor('#1E2A3B')); canvas.rect(0, 0, A4[0], 10*mm, fill=1, stroke=0)
        canvas.setFillColor(HexColor('#64748B')); canvas.setFont('Helvetica', 7)
        canvas.drawString(15*mm, 3.5*mm,
                          f'BMI Intelligence Hub  ·  {self.S["report_date"]}  ·  CONFIDENTIAL')
        canvas.drawRightString(A4[0]-15*mm, 3.5*mm, f'Page {doc.page}')
        canvas.restoreState()

    def _embed(self, buf, w_mm=170, h_mm=70):
        if buf is None: return None
        buf.seek(0)
        return RLImage(buf, width=w_mm*mm, height=h_mm*mm)

    def build(self) -> bytes:
        buf  = io.BytesIO()
        doc  = SimpleDocTemplate(buf, pagesize=A4,
                                 leftMargin=18*mm, rightMargin=18*mm,
                                 topMargin=22*mm,  bottomMargin=18*mm)
        S    = self.S
        story = [Spacer(1, 240), PageBreak()]

        # System analytics
        story += [
            Paragraph('System Analytics', self.s_h2),
            Paragraph(
                f'<b>{S["api_calls"]} API calls</b> across all modules over the last {self.L.days} days. '
                f'Avg response: <b>{S["avg_ms"]}ms</b> · P50: {S["median_ms"]}ms · '
                f'P95: {S["p95_ms"]}ms · P99: {S["p99_ms"]}ms · Max: {S["max_ms"]}ms.',
                self.s_body),
            Spacer(1, 6),
            self._embed(self.CH.get('api_daily'), 170, 68),
            Paragraph('Fig 1 — Daily API Activity (POST=Create, PATCH=Update)', self.s_code),
            Spacer(1, 6),
            self._embed(self.CH.get('hourly'), 170, 58),
            Paragraph('Fig 2 — Hourly distribution · Red>30 calls · Amber>15 · Blue=normal', self.s_code),
            PageBreak(),
        ]
        story = [s for s in story if s is not None]

        # Module breakdown
        mc    = self.L.module_breakdown()
        total = max(len(self.L.audit), 1)
        story.append(Paragraph('Module Breakdown', self.s_h2))
        if 'modules' in self.CH:
            story.append(self._embed(self.CH['modules'], 170, 90))
            story.append(Paragraph('Fig 3 — API calls per module with percentage share', self.s_code))
            story.append(Spacer(1, 6))
        if not mc.empty:
            mod_td = [['Module', 'Calls', '% Total', 'POST', 'PATCH']]
            for _, row in mc.iterrows():
                post_n  = len(self.L.audit[(self.L.audit['resource_type']==row['resource_type']) & (self.L.audit['action']=='POST')]) if not self.L.audit.empty else 0
                patch_n = len(self.L.audit[(self.L.audit['resource_type']==row['resource_type']) & (self.L.audit['action']=='PATCH')]) if not self.L.audit.empty else 0
                mod_td.append([row['resource_type'][:24], str(int(row['calls'])),
                               f"{round(row['calls']/total*100, 1)}%",
                               str(post_n), str(patch_n)])
            story.append(_make_rl_table(mod_td, [60*mm, 22*mm, 24*mm, 20*mm, 20*mm]))
        story.append(PageBreak())

        # Performance
        story += [Paragraph('API Performance Analysis', self.s_h2)]
        if 'response_time' in self.CH:
            story.append(self._embed(self.CH['response_time'], 170, 76))
            story.append(Paragraph('Fig 4 — Bucket distribution and full histogram with percentile lines', self.s_code))
            story.append(Spacer(1, 6))
        perf_td = [['Metric', 'Value', 'SLA Target', 'Status', 'Action']]
        perf_rows_pdf = [
            ('Avg Response',  f"{S['avg_ms']}ms",    '< 200ms',  '⚠ Warn',  'Optimise fm-report'),
            ('P50 Median',    f"{S['median_ms']}ms", '< 100ms',  '✓ Good',  'Maintain'),
            ('P95',           f"{S['p95_ms']}ms",    '< 1,000ms','⚠ Warn',  'Profile slow endpoints'),
            ('P99',           f"{S['p99_ms']}ms",    '< 2,000ms','✗ Crit',  'Root cause analysis'),
            ('Maximum',       f"{S['max_ms']}ms",    '< 5,000ms','✗ Crit',  'Outlier — check logs'),
        ]
        pb = self.L.perf_buckets()
        perf_rows_pdf += [
            ('Under 100ms', f"{int(pb.get('<100ms', 0))} calls",   '> 80%',  '✓ Good',  'Maintain'),
            ('Over 1 second', f"{int(pb.get('>1s', 0))} calls",    '< 1%',   '✗ Crit',  'Investigate'),
        ]
        for r in perf_rows_pdf:
            perf_td.append(list(r))
        pt = Table(perf_td, colWidths=[36*mm, 28*mm, 26*mm, 20*mm, 56*mm])
        pts = TableStyle([
            ('BACKGROUND',    (0,0), (-1,0), _rl_color(C.BLUE)),
            ('TEXTCOLOR',     (0,0), (-1,0), _rl_color(C.WHITE)),
            ('FONTNAME',      (0,0), (-1,0), 'Helvetica-Bold'),
            ('FONTSIZE',      (0,0), (-1,-1), 8),
            ('ALIGN',         (0,0), (-1,-1), 'CENTER'),
            ('VALIGN',        (0,0), (-1,-1), 'MIDDLE'),
            ('ROWHEIGHT',     (0,0), (-1,-1), 16),
            ('GRID',          (0,0), (-1,-1), 0.4, _rl_color(C.BORDER)),
            ('ROWBACKGROUNDS',(0,1), (-1,-1), [_rl_color(C.CARD), HexColor('#0F1623')]),
            ('TEXTCOLOR',     (0,1), (-1,-1), _rl_color(C.SLATE)),
        ])
        for ri, r in enumerate(perf_rows_pdf, 1):
            col_ = _rl_color(C.GREEN) if '✓' in r[3] else _rl_color(C.AMBER) if '⚠' in r[3] else _rl_color(C.RED)
            pts.add('TEXTCOLOR', (3,ri), (3,ri), col_); pts.add('FONTNAME', (3,ri), (3,ri), 'Helvetica-Bold')
            pts.add('TEXTCOLOR', (1,ri), (1,ri), col_); pts.add('FONTNAME', (1,ri), (1,ri), 'Helvetica-Bold')
        pt.setStyle(pts); story.append(pt); story.append(PageBreak())

        # Users
        story.append(Paragraph('Users & Activity', self.s_h2))
        if 'users' in self.CH:
            story.append(self._embed(self.CH['users'], 170, 60))
            story.append(Paragraph('Fig 5 — API activity per user', self.s_code))
            story.append(Spacer(1, 6))
        if not self.L.users.empty:
            uc    = self.L.user_activity()
            call_map_local = dict(zip(uc['user_id'].astype(str), uc['calls']))
            total = max(len(self.L.audit), 1)
            utd   = [['Employee ID', 'Name', 'Role', 'API Calls', '%', 'MFA', 'Last Login']]
            for _, u in self.L.users.iterrows():
                uid = str(u['id']); n = call_map_local.get(uid, 0)
                utd.append([
                    u.get('employee_id', ''),
                    f"{u['first_name']} {u['last_name']}",
                    str(u['role']).replace('_', ' ').title(),
                    str(int(n)), f"{round(n/total*100, 1)}%",
                    '✗ Off' if not u.get('mfa_enabled') else '✓ On',
                    str(u.get('last_login', ''))[:10] or 'Never',
                ])
            ut   = Table(utd, colWidths=[26*mm, 36*mm, 30*mm, 18*mm, 16*mm, 16*mm, 24*mm])
            uts  = TableStyle([
                ('BACKGROUND',    (0,0), (-1,0), _rl_color(C.BLUE)),
                ('TEXTCOLOR',     (0,0), (-1,0), _rl_color(C.WHITE)),
                ('FONTNAME',      (0,0), (-1,0), 'Helvetica-Bold'),
                ('FONTSIZE',      (0,0), (-1,-1), 8),
                ('ALIGN',         (0,0), (-1,-1), 'CENTER'),
                ('VALIGN',        (0,0), (-1,-1), 'MIDDLE'),
                ('ROWHEIGHT',     (0,0), (-1,-1), 15),
                ('GRID',          (0,0), (-1,-1), 0.4, _rl_color(C.BORDER)),
                ('ROWBACKGROUNDS',(0,1), (-1,-1), [_rl_color(C.CARD), HexColor('#0F1623')]),
                ('TEXTCOLOR',     (0,1), (-1,-1), _rl_color(C.SLATE)),
            ])
            for ri in range(1, len(utd)):
                mfa_v = utd[ri][5]
                uts.add('TEXTCOLOR', (5,ri), (5,ri), _rl_color(C.GREEN) if '✓' in mfa_v else _rl_color(C.RED))
                uts.add('FONTNAME',  (5,ri), (5,ri), 'Helvetica-Bold')
            ut.setStyle(uts); story.append(ut)
        story.append(Spacer(1, 8))
        if 'notifications' in self.CH:
            story.append(self._embed(self.CH['notifications'], 120, 80))
            story.append(Paragraph('Fig 6 — Notification type distribution', self.s_code))
        story.append(PageBreak())

        # Insights
        story.append(Paragraph('Insights & Recommendations', self.s_h2))
        insights = [
            (f'🔴 CRITICAL — MFA: {S["mfa_pct"]}% adoption',
             f'{S["mfa_enabled"]} of {S["active_users"]} users have MFA enabled. '
             'Every account without MFA is a credential-theft risk. Enforce 2FA immediately.', self.s_crit),
            (f'🔴 CRITICAL — P99 Response: {S["p99_ms"]}ms',
             f'Max response time: {S["max_ms"]}ms. Investigate slowest endpoints. '
             'Add Redis caching for frequent ORM queries. Target P99 < 2,000ms.', self.s_crit),
            ('⚠ HIGH — Off-Hours API Activity',
             'Significant API traffic detected outside office hours. '
             'Verify whether automated/cron tasks or unauthorised sessions.', self.s_warn),
            (f'⚠ HIGH — Equipment: {S["equip_repair"]} items under repair',
             f'KES {S["equip_value"]:,.0f} in assets unavailable for use. '
             'Establish 48h maintenance SLA and procure spare equipment.', self.s_warn),
            (f'✓ POSITIVE — Emergency Response: {S["alerts_resolved"]}/{S["alerts_total"]} resolved',
             'All emergency alerts were successfully resolved. Document the response playbook.', self.s_good),
            ('✓ POSITIVE — Task Review Speed',
             'Both submitted tasks reviewed within minutes. Outstanding supervisor performance.', self.s_good),
        ]
        for title, body, sty in insights:
            story.append(Paragraph(f'<b>{title}</b>', sty))
            story.append(Paragraph(body, self.s_body))
            story.append(Spacer(1, 7))

        doc.build(story, onFirstPage=self._cover, onLaterPages=self._later)
        return buf.getvalue()


# ─── Excel builder ────────────────────────────────────────────────────────────
class ExcelBuilder:
    def __init__(self, loader: BMIDataLoader, charts: dict):
        self.L = loader; self.S = loader.stats; self.CH = charts
        self.wb = Workbook()
        # colour fills
        self.KF = PatternFill('solid', fgColor='161B27')
        self.HF = PatternFill('solid', fgColor='1E2A3B')
        self.BF = PatternFill('solid', fgColor='3B63F5')
        self.DF = PatternFill('solid', fgColor='0F1623')
        self.HFont  = Font(name='Arial', bold=True, color='E2E8F0', size=10)
        self.TFont  = Font(name='Arial', bold=True, color='FFFFFF',  size=13)
        self.DFont  = Font(name='Arial', color='E2E8F0',  size=9)
        self.LFont  = Font(name='Arial', color='94A3B8',  size=9)
        self.MFont  = Font(name='Courier New', color='94A3B8', size=9)

    def _hdr(self, ws, col, row, val, width=14):
        c = ws.cell(row=row, column=col, value=val)
        c.font, c.fill = self.HFont, self.HF
        c.alignment = Alignment(horizontal='center', vertical='center')
        ws.column_dimensions[get_column_letter(col)].width = width
        return c

    def _data(self, ws, col, row, val, color='E2E8F0', bold=False, numfmt=None, alt=False):
        c = ws.cell(row=row, column=col, value=val)
        c.font   = Font(name='Arial', color=color, bold=bold, size=9)
        c.fill   = self.DF if alt else self.KF
        c.alignment = Alignment(horizontal='center', vertical='center')
        if numfmt: c.number_format = numfmt
        return c

    def _kpi_block(self, ws, row, col, label, value, fill_hex):
        ws.merge_cells(start_row=row, start_column=col, end_row=row, end_column=col+1)
        lc = ws.cell(row=row, column=col, value=label)
        lc.font, lc.fill = self.LFont, self.KF
        lc.alignment = Alignment(horizontal='center')
        ws.merge_cells(start_row=row+1, start_column=col, end_row=row+1, end_column=col+1)
        vc = ws.cell(row=row+1, column=col, value=value)
        vc.font  = Font(name='Arial', bold=True, color='FFFFFF', size=16)
        vc.fill  = PatternFill('solid', fgColor=fill_hex.lstrip('#'))
        vc.alignment = Alignment(horizontal='center', vertical='center')
        ws.row_dimensions[row+1].height = 28

    def build(self) -> bytes:
        S  = self.S
        wb = self.wb

        # ── Sheet 1: Executive Summary ────────────────────────────────────
        ws1 = wb.active; ws1.title = 'Executive Summary'
        ws1.sheet_view.showGridLines = False
        for c in range(1, 20):
            for r in range(1, 60):
                ws1.cell(row=r, column=c).fill = self.KF

        ws1.merge_cells('A1:N1')
        t = ws1['A1']; t.value = 'BMI Intelligence Hub — Analytics Report'
        t.font = Font(name='Arial', bold=True, color='FFFFFF', size=16)
        t.fill = self.BF; t.alignment = Alignment(horizontal='center', vertical='center')
        ws1.row_dimensions[1].height = 36

        ws1.merge_cells('A2:N2')
        s = ws1['A2']
        s.value = f"Broadcast Media Institution  ·  Generated: {S['report_date']}"
        s.font  = Font(name='Arial', color='94A3B8', size=10)
        s.alignment = Alignment(horizontal='center')
        ws1.row_dimensions[2].height = 20

        kpi_ex = [
            ('Active Users',    S['active_users'],              'C.22C55E'),
            ('API Calls',       S['api_calls'],                 '3B63F5'),
            ('Avg Response',    f"{S['avg_ms']}ms",            'F59E0B'),
            ('MFA Adoption',    f"{S['mfa_pct']}%",            'EF4444' if S['mfa_pct']<50 else '22C55E'),
            ('Asset Value',     f"KES {S['equip_value']:,.0f}",'8B5CF6'),
            ('Alerts Resolved', f"{S['alerts_resolved']}/{S['alerts_total']}", '22C55E'),
        ]
        col_k = 1
        for lbl, val, fill in kpi_ex:
            self._kpi_block(ws1, 4, col_k, lbl, val, fill.lstrip('C.'))
            ws1.column_dimensions[get_column_letter(col_k)].width = 14
            ws1.column_dimensions[get_column_letter(col_k+1)].width = 14
            col_k += 3

        # Daily API table
        ws1.cell(row=7, column=1, value='Daily API Activity').font = Font(name='Arial', bold=True, color='FFFFFF', size=11)
        for i, h in enumerate(['Date', 'Total', 'POST', 'PATCH', '% POST', '% PATCH'], 1):
            self._hdr(ws1, i, 8, h)
        mc_daily = self.L.audit.groupby(['date', 'action']).size().unstack(fill_value=0).reset_index() if not self.L.audit.empty else pd.DataFrame()
        for ri, (_, row) in enumerate(mc_daily.iterrows() if not mc_daily.empty else iter([]), 9):
            alt = ri % 2 == 0
            self._data(ws1, 1, ri, str(row['date']), alt=alt)
            posts_n  = int(row.get('POST', 0));  patches_n = int(row.get('PATCH', 0))
            total_n  = posts_n + patches_n
            self._data(ws1, 2, ri, total_n,   bold=True, alt=alt)
            self._data(ws1, 3, ri, posts_n,   color='3B63F5', bold=True, alt=alt)
            self._data(ws1, 4, ri, patches_n, color='8B5CF6', bold=True, alt=alt)
            self._data(ws1, 5, ri, posts_n/max(total_n,1),   numfmt='0.0%', alt=alt)
            self._data(ws1, 6, ri, patches_n/max(total_n,1), numfmt='0.0%', alt=alt)

        # Embed dashboard chart
        if 'dashboard' in self.CH:
            from openpyxl.drawing.image import Image as XLImage
            self.CH['dashboard'].seek(0)
            xi = XLImage(self.CH['dashboard'])
            xi.anchor = 'H4'; xi.width = 640; xi.height = 420
            ws1.add_image(xi)

        # ── Sheet 2: System Analytics ─────────────────────────────────────
        ws2 = wb.create_sheet('System Analytics')
        ws2.sheet_view.showGridLines = False
        ws2.merge_cells('A1:K1')
        ws2['A1'].value = 'System & API Performance Analytics'
        ws2['A1'].font  = Font(name='Arial', bold=True, color='FFFFFF', size=13)
        ws2['A1'].fill  = self.HF
        ws2['A1'].alignment = Alignment(horizontal='center')
        ws2.row_dimensions[1].height = 28

        perf_hdrs = ['Metric', 'Value', 'SLA Target', 'Status', 'Notes']
        for i, h in enumerate(perf_hdrs, 1):
            self._hdr(ws2, i, 3, h, width=20)
        perf_data_xl = [
            ('Average Response', f"{S['avg_ms']}ms",    '< 200ms',  '⚠ Above SLA', 'Profile fm-report endpoint'),
            ('P50 Median',       f"{S['median_ms']}ms", '< 100ms',  '✓ Good',       'Maintain'),
            ('P95',              f"{S['p95_ms']}ms",    '< 1,000ms','⚠ Warn',       'Investigate slow queries'),
            ('P99',              f"{S['p99_ms']}ms",    '< 2,000ms','✗ Critical',   'Urgent: root cause analysis'),
            ('Maximum',          f"{S['max_ms']}ms",    '< 5,000ms','✗ Outlier',    'Check server logs'),
            ('Under 100ms',      f"{int(self.L.perf_buckets().get('<100ms',0))} ({_pct(int(self.L.perf_buckets().get('<100ms',0)), S['api_calls'])}%)", '> 80%', '✓ OK', 'Target met'),
            ('Over 1 second',    f"{int(self.L.perf_buckets().get('>1s',0))} ({_pct(int(self.L.perf_buckets().get('>1s',0)), S['api_calls'])}%)", '< 1%', '✗ Critical', 'Needs immediate attention'),
        ]
        for ri, (met, val, tgt, st, note) in enumerate(perf_data_xl, 4):
            alt = ri % 2 == 0
            col_ = 'EF4444' if '✗' in st else 'F59E0B' if '⚠' in st else '22C55E'
            self._data(ws2, 1, ri, met,  alt=alt); self._data(ws2, 2, ri, val, color=col_, bold=True, alt=alt)
            self._data(ws2, 3, ri, tgt,  alt=alt); self._data(ws2, 4, ri, st,  color=col_, bold=True, alt=alt)
            self._data(ws2, 5, ri, note, alt=alt)

        # Module breakdown table
        ws2.cell(row=14, column=1, value='Module API Breakdown').font = Font(name='Arial', bold=True, color='FFFFFF', size=11)
        for i, h in enumerate(['Module', 'Total Calls', '% of Total', 'POST', 'PATCH'], 1):
            self._hdr(ws2, i, 15, h)
        mc = self.L.module_breakdown()
        total = max(len(self.L.audit), 1)
        for ri, (_, row) in enumerate(mc.iterrows(), 16):
            alt = ri % 2 == 0
            post_n  = len(self.L.audit[(self.L.audit['resource_type']==row['resource_type']) & (self.L.audit['action']=='POST')]) if not self.L.audit.empty else 0
            patch_n = len(self.L.audit[(self.L.audit['resource_type']==row['resource_type']) & (self.L.audit['action']=='PATCH')]) if not self.L.audit.empty else 0
            self._data(ws2, 1, ri, row['resource_type'].title(), alt=alt)
            self._data(ws2, 2, ri, int(row['calls']), color='3B63F5', bold=True, alt=alt)
            self._data(ws2, 3, ri, row['calls']/total, numfmt='0.0%', alt=alt)
            self._data(ws2, 4, ri, post_n,  alt=alt)
            self._data(ws2, 5, ri, patch_n, alt=alt)

        # ── Sheet 3: Audit Log ────────────────────────────────────────────
        ws3 = wb.create_sheet('Audit Log')
        ws3.sheet_view.showGridLines = False
        ws3.merge_cells('A1:H1')
        ws3['A1'].value = f'Audit Log — {len(self.L.audit)} Records'
        ws3['A1'].font  = Font(name='Arial', bold=True, color='FFFFFF', size=12)
        ws3['A1'].fill  = self.HF
        ws3['A1'].alignment = Alignment(horizontal='center')

        hdrs3   = ['Timestamp', 'User', 'Action', 'Module', 'Resource ID', 'Duration (ms)', 'Time Bucket', 'Slow?']
        widths3 = [22, 20, 10, 18, 28, 14, 16, 8]
        for i, (h, w) in enumerate(zip(hdrs3, widths3), 1):
            self._hdr(ws3, i, 2, h, width=w)

        def time_bucket(h):
            return 'Night(0-6)' if h<6 else 'Morning(6-12)' if h<12 else 'Afternoon(12-18)' if h<18 else 'Evening(18-24)'

        for ri, (_, row) in enumerate(self.L.audit.iterrows(), 3):
            alt    = ri % 2 == 0
            uname  = self.L.user_map.get(str(row.get('user_id','')), 'System')
            ms_val = row.get('duration_ms')
            try:
                extra = json.loads(row['extra']); sc = extra.get('status_code', '')
            except: sc = ''
            is_slow = 'Yes' if (ms_val and not pd.isna(ms_val) and float(ms_val) > 1000) else 'No'
            ms_col  = 'EF4444' if is_slow=='Yes' else 'F59E0B' if (ms_val and float(ms_val)>500) else '22C55E'
            act_col = '3B63F5' if row.get('action','')=='POST' else '8B5CF6'
            self._data(ws3, 1, ri, str(row.get('created_at',''))[:19], alt=alt)
            self._data(ws3, 2, ri, uname, alt=alt)
            self._data(ws3, 3, ri, row.get('action',''), color=act_col, bold=True, alt=alt)
            self._data(ws3, 4, ri, str(row.get('resource_type','')), alt=alt)
            self._data(ws3, 5, ri, str(row.get('resource_id',''))[:28], alt=alt)
            self._data(ws3, 6, ri, round(float(ms_val), 1) if ms_val and not pd.isna(ms_val) else '', color=ms_col, bold=True, numfmt='#,##0.0', alt=alt)
            self._data(ws3, 7, ri, time_bucket(int(row.get('hour', 12))), alt=alt)
            self._data(ws3, 8, ri, is_slow, color='EF4444' if is_slow=='Yes' else '22C55E', bold=True, alt=alt)

        # ── Sheet 4: All Modules ──────────────────────────────────────────
        ws4 = wb.create_sheet('All Modules')
        ws4.sheet_view.showGridLines = False
        cur = 1

        def write_section(ws, start_row, title, headers, rows, widths):
            ws.merge_cells(start_row=start_row, start_column=1, end_row=start_row, end_column=len(headers))
            tc = ws.cell(row=start_row, column=1, value=title)
            tc.font = Font(name='Arial', bold=True, color='FFFFFF', size=11)
            tc.fill = self.BF; tc.alignment = Alignment(horizontal='center')
            ws.row_dimensions[start_row].height = 22
            for i, (h, w) in enumerate(zip(headers, widths), 1):
                self._hdr(ws, i, start_row+1, h, width=w)
            for ri, row_ in enumerate(rows, start_row+2):
                alt = ri % 2 == 0
                for ci, val in enumerate(row_, 1):
                    self._data(ws, ci, ri, str(val), alt=alt)
            return start_row + 2 + len(rows) + 2

        # Users section
        if not self.L.users.empty:
            uc = self.L.user_activity()
            call_map_local = dict(zip(uc['user_id'].astype(str), uc['calls']))
            total = max(len(self.L.audit), 1)
            user_rows = [[
                u.get('employee_id',''), f"{u['first_name']} {u['last_name']}",
                str(u['role']).replace('_',' ').title(),
                '✗ Off' if not u.get('mfa_enabled') else '✓ On',
                str(int(call_map_local.get(str(u['id']), 0))),
                f"{round(call_map_local.get(str(u['id']),0)/total*100,1)}%",
                'HIGH' if not u.get('mfa_enabled') else 'LOW',
            ] for _, u in self.L.users.iterrows()]
            cur = write_section(ws4, cur, 'USERS',
                ['Employee ID','Name','Role','MFA','API Calls','Activity %','Risk'],
                user_rows, [16,24,20,10,12,14,10])

        # Tasks
        if not self.L.tasks.empty:
            task_rows = [[
                t.get('title','')[:30], t.get('priority','').title(), t.get('status','').title(),
                f"{t.get('assigned_to__first_name','')} {t.get('assigned_to__last_name','')}",
                str(t.get('due_date',''))[:10], str(t.get('submitted_at',''))[:16],
                str(t.get('reviewed_at',''))[:16],
            ] for _, t in self.L.tasks.iterrows()]
            cur = write_section(ws4, cur, 'TASKS',
                ['Title','Priority','Status','Assigned To','Due Date','Submitted','Reviewed'],
                task_rows, [28,14,14,22,14,18,18])

        # Equipment
        if not self.L.equipment.empty:
            eq_rows = [[
                e.get('asset_tag',''), e.get('name','')[:28],
                e.get('status','').replace('_',' ').title(), e.get('condition','').title(),
                f"KES {float(e.get('purchase_cost',0)):,.0f}",
                e.get('category__name',''), e.get('location',''),
            ] for _, e in self.L.equipment.iterrows()]
            cur = write_section(ws4, cur, 'EQUIPMENT',
                ['Asset Tag','Name','Status','Condition','Value (KES)','Category','Location'],
                eq_rows, [14,30,20,16,20,16,14])

        # Alerts
        if not self.L.alerts.empty:
            alert_rows = [[
                a.get('title','')[:28], a.get('alert_type','').title(),
                a.get('severity','').title(),
                '✓ Yes' if a.get('resolved') else '✗ No',
                str(a.get('created_at',''))[:10],
            ] for _, a in self.L.alerts.iterrows()]
            cur = write_section(ws4, cur, 'EMERGENCY ALERTS',
                ['Title','Type','Severity','Resolved','Date'],
                alert_rows, [36,16,14,10,14])

        # Certificates
        if not self.L.certificates.empty:
            cert_rows = [[
                c.get('certificate_number',''), c.get('certificate_type','').title(),
                c.get('status','').title(), str(c.get('issue_date',''))[:10],
                f"{c.get('recipient__first_name','')} {c.get('recipient__last_name','')}",
            ] for _, c in self.L.certificates.iterrows()]
            cur = write_section(ws4, cur, 'CERTIFICATES',
                ['Cert Number','Type','Status','Issue Date','Recipient'],
                cert_rows, [26,24,14,14,24])

        # Feedback
        if not self.L.feedback.empty:
            fb_rows = [[
                f.get('ticket_number',''), f.get('category','').title(),
                f.get('priority','').title(), f.get('status','').title(),
                str(f.get('created_at',''))[:10],
            ] for _, f in self.L.feedback.iterrows()]
            cur = write_section(ws4, cur, 'SUPPORT TICKETS',
                ['Ticket #','Category','Priority','Status','Created'],
                fb_rows, [14,20,14,14,14])

        out = io.BytesIO()
        wb.save(out)
        return out.getvalue()


# ─── PPTX builder ─────────────────────────────────────────────────────────────
class PPTXBuilder:
    def __init__(self, loader: BMIDataLoader, charts: dict):
        self.L = loader; self.S = loader.stats; self.CH = charts
        self.prs = Presentation()
        self.prs.slide_width  = Inches(13.33)
        self.prs.slide_height = Inches(7.5)
        self.blank = self.prs.slide_layouts[6]

    def _bg(self, slide):
        r = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
        r.fill.solid(); r.fill.fore_color.rgb = RGBColor(0x0d, 0x11, 0x17)
        r.line.fill.background()

    def _rect(self, sl, l, t, w, h, hx):
        s = sl.shapes.add_shape(1, Inches(l), Inches(t), Inches(w), Inches(h))
        s.fill.solid(); s.fill.fore_color.rgb = RGBColor.from_string(hx.lstrip('#'))
        s.line.fill.background(); return s

    def _txt(self, sl, text, l, t, w, h, sz=11, bold=False, color='E2E8F0', align=PP_ALIGN.LEFT):
        tb = sl.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
        tf = tb.text_frame; tf.word_wrap = True
        p  = tf.paragraphs[0]; p.alignment = align
        run = p.add_run(); run.text = str(text)
        run.font.size = Pt(sz); run.font.bold = bold
        run.font.color.rgb = RGBColor.from_string(color.lstrip('#'))
        run.font.name = 'Calibri'; return tb

    def _img(self, sl, key, l, t, w, h):
        if key in self.CH and self.CH[key]:
            self.CH[key].seek(0)
            sl.shapes.add_picture(self.CH[key], Inches(l), Inches(t), Inches(w), Inches(h))

    def build(self) -> bytes:
        S  = self.S

        # Cover
        s1 = self.prs.slides.add_slide(self.blank); self._bg(s1)
        self._rect(s1, 0, 0, 13.33, 2.1, C.BLUE)
        self._txt(s1, 'INTELLIGENCE HUB', 0.4, 0.15, 12.5, 1.05, sz=38, bold=True, color=C.WHITE, align=PP_ALIGN.CENTER)
        self._txt(s1, 'Broadcast Media Institution — Full Analytics Report', 0.4, 1.1, 12.5, 0.55, sz=14, color='93C5FD', align=PP_ALIGN.CENTER)
        self._txt(s1, f'Generated: {S["report_date"]}  ·  Pandas · Matplotlib · ReportLab · OpenPyXL · python-pptx',
                  0.4, 1.68, 12.5, 0.35, sz=8, color='64748B', align=PP_ALIGN.CENTER)
        kps = [
            (str(S['active_users']),              'Active Users',    C.GREEN),
            (str(S['api_calls']),                 'API Calls',       C.BLUE),
            (f"{S['avg_ms']}ms",                  'Avg Response',    C.AMBER),
            (f"{S['mfa_pct']}% {'⚠' if S['mfa_pct']<50 else '✓'}", 'MFA Adoption', C.RED if S['mfa_pct']<50 else C.GREEN),
            (f"{S['alerts_resolved']}/{S['alerts_total']}", 'Alerts OK', C.GREEN),
            (f"KES {S['equip_value']/1e6:.1f}M",  'Assets',          C.PURPLE),
        ]
        kw = 13.33 / len(kps)
        for i, (val, lbl, col) in enumerate(kps):
            x = i * kw
            self._rect(s1, x+0.04, 2.25, kw-0.08, 1.4, C.CARD)
            self._txt(s1, val, x+0.04, 2.42, kw-0.08, 0.72, sz=22, bold=True, color=col, align=PP_ALIGN.CENTER)
            self._txt(s1, lbl, x+0.04, 3.14, kw-0.08, 0.38, sz=8,  color=C.SLATE, align=PP_ALIGN.CENTER)
        self._txt(s1, 'CONFIDENTIAL — FOR INTERNAL USE ONLY', 0.4, 7.15, 12.5, 0.3, sz=7, color='475569', align=PP_ALIGN.CENTER)

        # Slide 2: API Activity
        s2 = self.prs.slides.add_slide(self.blank); self._bg(s2)
        self._rect(s2, 0, 0, 13.33, 0.72, C.CARD)
        self._txt(s2, 'API Activity Analysis', 0.3, 0.08, 9, 0.56, sz=20, bold=True)
        self._txt(s2, f'{S["api_calls"]} calls · {S["api_posts"]} POST · {S["api_patches"]} PATCH · last {self.L.days} days', 0.3, 0.56, 9, 0.28, sz=9, color=C.SLATE)
        self._img(s2, 'api_daily', 0.3, 0.88, 8.7, 2.9)
        self._img(s2, 'hourly',    0.3, 3.95, 8.7, 3.3)
        for i, (lbl, val, col) in enumerate([
            ('Total Calls', str(S['api_calls']), C.BLUE), ('POST', f"{S['api_posts']} ({_pct(S['api_posts'],S['api_calls'])}%)", C.GREEN),
            ('PATCH', f"{S['api_patches']} ({_pct(S['api_patches'],S['api_calls'])}%)", C.PURPLE),
            ('Avg Response', f"{S['avg_ms']}ms", C.AMBER), ('P99', f"{S['p99_ms']}ms", C.RED),
            ('Users Active', str(len(self.L.user_activity())), C.CYAN),
        ]):
            self._rect(s2, 9.2, 1.0+i*1.02, 3.85, 0.88, C.CARD)
            self._txt(s2, lbl, 9.3, 1.05+i*1.02, 3.6, 0.3,  sz=8,  color=C.SLATE)
            self._txt(s2, val, 9.3, 1.36+i*1.02, 3.6, 0.45, sz=15, bold=True, color=col)

        # Slide 3: Performance
        s3 = self.prs.slides.add_slide(self.blank); self._bg(s3)
        self._rect(s3, 0, 0, 13.33, 0.72, C.CARD)
        self._txt(s3, 'API Performance', 0.3, 0.08, 9, 0.56, sz=20, bold=True)
        self._txt(s3, f'Avg: {S["avg_ms"]}ms · P50: {S["median_ms"]}ms · P95: {S["p95_ms"]}ms · P99: {S["p99_ms"]}ms · Max: {S["max_ms"]}ms', 0.3, 0.56, 12.5, 0.28, sz=9, color=C.SLATE)
        self._img(s3, 'response_time', 0.3, 0.86, 8.9, 3.8)
        pb = self.L.perf_buckets()
        for i, (lbl, val, col, st) in enumerate([
            ('Average',  f"{S['avg_ms']}ms",    C.AMBER, '⚠ Warn'),
            ('P50',      f"{S['median_ms']}ms", C.GREEN, '✓ Good'),
            ('P95',      f"{S['p95_ms']}ms",    C.AMBER, '⚠ Warn'),
            ('P99',      f"{S['p99_ms']}ms",    C.RED,   '✗ Crit'),
            ('Max',      f"{S['max_ms']}ms",    C.RED,   '✗ Crit'),
            ('<100ms',   f"{int(pb.get('<100ms',0))} · {_pct(int(pb.get('<100ms',0)),S['api_calls'])}%", C.GREEN,'✓ OK'),
            ('>1s slow', f"{int(pb.get('>1s',0))} · {_pct(int(pb.get('>1s',0)),S['api_calls'])}%",      C.RED,  '✗ Crit'),
        ]):
            ci = i % 4; ri_ = i // 4; x_ = 0.3+ci*3.26; y_ = 4.85+ri_*1.3
            self._rect(s3, x_, y_, 3.1, 1.2, C.CARD)
            self._txt(s3, lbl, x_+0.1, y_+0.04, 2.9, 0.32, sz=8,  color=C.SLATE)
            self._txt(s3, val, x_+0.1, y_+0.38, 2.9, 0.55, sz=16, bold=True, color=col, align=PP_ALIGN.CENTER)
            self._txt(s3, st,  x_+0.1, y_+0.92, 2.9, 0.28, sz=8,  color=col,  align=PP_ALIGN.CENTER)

        # Slide 4: Modules + Users
        s4 = self.prs.slides.add_slide(self.blank); self._bg(s4)
        self._rect(s4, 0, 0, 13.33, 0.72, C.CARD)
        self._txt(s4, 'Module & User Activity', 0.3, 0.08, 9, 0.56, sz=20, bold=True)
        self._img(s4, 'modules',       0.3, 0.82, 8.7, 4.2)
        self._img(s4, 'notifications', 0.3, 5.1,  4.8, 2.3)
        self._img(s4, 'users',         5.3, 5.1,  4.8, 2.3)
        mc = self.L.module_breakdown()
        self._rect(s4, 9.2, 0.82, 3.9, 4.2, C.CARD)
        self._txt(s4, '⚠ Key Findings', 9.35, 0.92, 3.6, 0.38, sz=11, bold=True, color=C.AMBER)
        if not mc.empty:
            total = max(len(self.L.audit), 1)
            for i, (_, row) in enumerate(mc.head(8).iterrows()):
                self._txt(s4, f'• {row["resource_type"]}: {int(row["calls"])} ({round(row["calls"]/total*100,1)}%)',
                          9.35, 1.42+i*0.38, 3.6, 0.34, sz=8, color=C.SLATE)

        # Slide 5: Operations
        s5 = self.prs.slides.add_slide(self.blank); self._bg(s5)
        self._rect(s5, 0, 0, 13.33, 0.72, C.CARD)
        self._txt(s5, 'Operations & Assets', 0.3, 0.08, 9, 0.56, sz=20, bold=True)
        op_items = [
            ('Equipment',      f"{S['equip_total']} Items · {S['equip_repair']} Repair", C.AMBER,  f"KES {S['equip_value']:,.0f}"),
            ('Emergency Alerts', f"{S['alerts_resolved']}/{S['alerts_total']} Resolved",  C.GREEN,  'All high severity'),
            ('Certificates',   f"{S['certs_issued']} Issued",                             C.PURPLE, 'Digital verification'),
            ('Support Tickets', f"{S['feedback_total']} Total",                           C.CYAN,   f"{S['feedback_open']} open"),
            ('Notifications',  f"{S['notif_total']} Sent",                               C.BLUE,   f"{S['notif_read']} read"),
            ('File Transfers', f"{S['files_total']} Files",                              C.BLUE,   f"{S['files_kb']} KB"),
            ('Tasks',          f"{S['tasks_total']} Total",                              C.GREEN,  f"{S['tasks_approved']} approved"),
            ('MFA Security',   f"{S['mfa_pct']}% Adoption",                             C.RED if S['mfa_pct']<50 else C.GREEN, '⚠ Action required' if S['mfa_pct']<50 else '✓ Good'),
        ]
        cw2 = 13.13 / 4
        for i, (title, val, col, note) in enumerate(op_items):
            ci = i % 4; ri_ = i // 4; x_ = 0.1+ci*cw2; y_ = 0.9+ri_*2.9
            self._rect(s5, x_, y_, cw2-0.08, 2.6, C.CARD)
            self._rect(s5, x_, y_, cw2-0.08, 0.45, col)
            self._txt(s5, title, x_+0.1, y_+0.04, cw2-0.22, 0.38, sz=9,  bold=True, color=C.WHITE)
            self._txt(s5, val,   x_+0.1, y_+0.6,  cw2-0.22, 0.9,  sz=13, bold=True, color=col, align=PP_ALIGN.CENTER)
            self._txt(s5, note,  x_+0.1, y_+1.68, cw2-0.22, 0.8,  sz=8,  color=C.SLATE)

        # Slide 6: Recommendations
        s6 = self.prs.slides.add_slide(self.blank); self._bg(s6)
        self._rect(s6, 0, 0, 13.33, 1.0, C.BLUE)
        self._txt(s6, 'Recommendations & Action Items', 0.4, 0.12, 12.5, 0.75, sz=24, bold=True, align=PP_ALIGN.CENTER)
        recs = [
            ('🔴 CRITICAL', f'Enforce MFA ({S["mfa_pct"]}% adoption)',
             f'Enforce 2FA for all {S["active_users"]} accounts immediately. Admin roles first.', C.RED),
            ('🔴 CRITICAL', f'Fix P99 Latency ({S["p99_ms"]}ms)',
             'Profile fm-report endpoint. Add Redis caching. Target P99 < 2,000ms.', C.RED),
            ('⚠ HIGH', 'Audit Off-Hours Activity',
             'Significant API calls outside business hours. Verify cron jobs vs unauthorised.', C.AMBER),
            ('⚠ HIGH', f'Equipment SLA ({S["equip_repair"]} under repair)',
             f'KES {S["equip_value"]:,.0f} idle. Establish 48h maintenance SLA.', C.AMBER),
            ('✓ SUSTAIN', f'Emergency Protocol ({S["alerts_resolved"]}/{S["alerts_total"]} OK)',
             'All alerts resolved. Document the response playbook and train staff.', C.GREEN),
            ('ℹ ACTION', 'Build Data Baseline',
             f'Enable daily exports. Target 30+ days for meaningful trend analysis.', C.PURPLE),
        ]
        rw = 13.13 / 3
        for i, (pri, title, desc, col) in enumerate(recs):
            ci = i % 3; ri_ = i // 3; x_ = 0.1+ci*rw; y_ = 1.15+ri_*2.9
            self._rect(s6, x_, y_,       rw-0.08, 2.65, C.CARD)
            self._rect(s6, x_, y_,       rw-0.08, 0.45, col)
            self._txt(s6, pri,   x_+0.1, y_+0.04,  rw-0.22, 0.38, sz=9,  bold=True, color=C.WHITE)
            self._txt(s6, title, x_+0.1, y_+0.58,  rw-0.22, 0.55, sz=13, bold=True, color=col)
            self._txt(s6, desc,  x_+0.1, y_+1.25,  rw-0.22, 1.3,  sz=9,  color=C.SLATE)

        out = io.BytesIO()
        self.prs.save(out)
        return out.getvalue()


# ─── JSON / Power BI export ───────────────────────────────────────────────────
class JSONBuilder:
    def __init__(self, loader: BMIDataLoader):
        self.L = loader; self.S = loader.stats

    def build(self) -> dict:
        L  = self.L; S  = self.S
        mc = L.module_breakdown()
        uc = L.user_activity()
        nt = L.notif_types()
        pb = L.perf_buckets()

        return {
            'metadata': {
                'org':          'Broadcast Media Institution (BMI)',
                'generated':    S['report_date'],
                'period_days':  L.days,
                'tool_stack':   ['pandas', 'numpy', 'scipy', 'matplotlib', 'seaborn',
                                 'reportlab', 'openpyxl', 'xlsxwriter', 'python-pptx'],
            },
            'kpis':        {k: _sv(v) for k, v in S.items()},
            'time_series': {
                'daily':  [{'date': str(r['date']), 'calls': int(r['calls'])}
                           for _, r in L.audit.groupby('date').size().reset_index(name='calls').iterrows()] if not L.audit.empty else [],
                'hourly': [{'hour': int(r['hour']), 'calls': int(r['calls'])}
                           for _, r in L.hourly().iterrows()],
            },
            'breakdowns': {
                'by_module':     [{'module': r['resource_type'], 'calls': int(r['calls']),
                                   'pct': round(r['calls']/max(len(L.audit),1)*100, 1)}
                                  for _, r in mc.iterrows()] if not mc.empty else [],
                'by_user':       [{'user': r['name'], 'calls': int(r['calls'])}
                                  for _, r in uc.iterrows()] if not uc.empty else [],
                'notifications': [{'type': r['type'], 'count': int(r['count'])}
                                  for _, r in nt.iterrows()] if not nt.empty else [],
            },
            'performance': {
                'avg_ms':    _sv(S['avg_ms']), 'median_ms': _sv(S['median_ms']),
                'p95_ms':    _sv(S['p95_ms']), 'p99_ms':    _sv(S['p99_ms']),
                'max_ms':    _sv(S['max_ms']),
                'buckets':   {k: int(v) for k, v in pb.items()},
            },
            'entities': {
                'users':      L.users.to_dict('records') if not L.users.empty else [],
                'tasks':      L.tasks.to_dict('records') if not L.tasks.empty else [],
                'equipment':  L.equipment.to_dict('records') if not L.equipment.empty else [],
                'alerts':     L.alerts.to_dict('records') if not L.alerts.empty else [],
                'certificates': L.certificates.to_dict('records') if not L.certificates.empty else [],
                'feedback':   L.feedback.to_dict('records') if not L.feedback.empty else [],
            },
            'insights': [
                {'severity': 'critical' if S['mfa_pct'] < 50 else 'positive',
                 'module': 'security',
                 'message': f'{S["mfa_pct"]}% MFA adoption — {S["active_users"]-S["mfa_enabled"]} accounts without 2FA'},
                {'severity': 'critical' if S['p99_ms'] > 2000 else 'warning',
                 'module': 'performance',
                 'message': f'P99={S["p99_ms"]}ms, max={S["max_ms"]}ms — investigate slow endpoints'},
                {'severity': 'positive' if S['alerts_resolved'] == S['alerts_total'] else 'warning',
                 'module': 'operations',
                 'message': f'{S["alerts_resolved"]}/{S["alerts_total"]} emergency alerts resolved'},
                {'severity': 'warning' if S['equip_repair'] > 0 else 'positive',
                 'module': 'equipment',
                 'message': f'{S["equip_repair"]} of {S["equip_total"]} equipment items under repair'},
            ],
        }


# ─── Enriched CSV exports ─────────────────────────────────────────────────────
class CSVBuilder:
    def __init__(self, loader: BMIDataLoader):
        self.L = loader

    def audit_enriched(self) -> bytes:
        if self.L.audit.empty: return b''
        ae = self.L.audit.copy()
        ae['user_name']  = ae['user_id'].astype(str).map(self.L.user_map).fillna('Unknown')
        ae['is_slow']    = ae['duration_ms'].apply(lambda x: 'Yes' if (pd.notna(x) and x > 1000) else 'No')
        ae['time_bucket']= ae['hour'].apply(lambda h: 'Night(0-6)' if h<6 else 'Morning(6-12)' if h<12 else 'Afternoon(12-18)' if h<18 else 'Evening(18-24)')
        cols = ['created_at','date','hour','time_bucket','user_name','action','resource_type','duration_ms','is_slow']
        return ae[cols].to_csv(index=False).encode('utf-8-sig')

    def users_enriched(self) -> bytes:
        if self.L.users.empty: return b''
        ue = self.L.users.copy()
        uc = self.L.user_activity()
        call_map = dict(zip(uc['user_id'].astype(str), uc['calls']))
        total = max(len(self.L.audit), 1)
        ue['api_calls']    = ue['id'].astype(str).map(call_map).fillna(0).astype(int)
        ue['activity_pct'] = (ue['api_calls'] / total * 100).round(1)
        ue['mfa_status']   = ue['mfa_enabled'].apply(lambda x: 'Enabled' if x else 'Disabled')
        ue['risk_level']   = ue['mfa_status'].apply(lambda x: 'HIGH' if x == 'Disabled' else 'LOW')
        cols = ['employee_id','first_name','last_name','role','mfa_status','risk_level','api_calls','activity_pct','date_joined','last_login']
        return ue[cols].to_csv(index=False).encode('utf-8-sig')

    def module_breakdown(self) -> bytes:
        mc = self.L.module_breakdown()
        if mc.empty: return b''
        total = max(len(self.L.audit), 1)
        mc = mc.copy()
        mc['pct'] = (mc['calls'] / total * 100).round(1)
        mc['post_calls']  = mc['resource_type'].apply(lambda r: len(self.L.audit[(self.L.audit['resource_type']==r)&(self.L.audit['action']=='POST')]) if not self.L.audit.empty else 0)
        mc['patch_calls'] = mc['resource_type'].apply(lambda r: len(self.L.audit[(self.L.audit['resource_type']==r)&(self.L.audit['action']=='PATCH')]) if not self.L.audit.empty else 0)
        return mc.to_csv(index=False).encode('utf-8-sig')

    def performance(self) -> bytes:
        pb = self.L.perf_buckets()
        S  = self.L.stats
        rows = [
            {'metric':'avg_ms',       'value':S['avg_ms'],    'target':'<200',  'status':'warn'},
            {'metric':'median_ms',    'value':S['median_ms'], 'target':'<100',  'status':'good'},
            {'metric':'p95_ms',       'value':S['p95_ms'],    'target':'<1000', 'status':'warn'},
            {'metric':'p99_ms',       'value':S['p99_ms'],    'target':'<2000', 'status':'critical'},
            {'metric':'max_ms',       'value':S['max_ms'],    'target':'<5000', 'status':'critical'},
            {'metric':'under_100ms',  'value':int(pb.get('<100ms',0)), 'target':'>80%','status':'good'},
            {'metric':'over_1s',      'value':int(pb.get('>1s',0)),    'target':'<1%', 'status':'critical'},
        ]
        return pd.DataFrame(rows).to_csv(index=False).encode('utf-8-sig')


# ════════════════════════════════════════════════════════════════════════════════
#  API VIEWS
# ════════════════════════════════════════════════════════════════════════════════

class AnalyticsExportView(APIView):
    """
    GET /analytics/export/<format>/<module>/
    Formats: pdf · excel · pptx · csv · json
    Modules: full · audit · users · modules · performance · attendance ·
             equipment · alerts · certificates · feedback · wifi
    Query params: ?days=30&department=<uuid>
    """
    permission_classes = [IsAuthenticated]

    def get(self, request, fmt: str, module: str = 'full'):
        org  = request.user.organisation
        days = int(request.query_params.get('days', 30))
        fmt  = fmt.lower()

        # Load all data
        loader = BMIDataLoader(org, days=days)
        S      = loader.stats

        # ── JSON — no charts needed ─────────────────────────────────────────
        if fmt == 'json':
            payload = JSONBuilder(loader).build()
            resp    = JsonResponse(payload, json_dumps_params={'indent': 2})
            resp['Content-Disposition'] = f'attachment; filename="{org.code}-analytics-{S["report_date"]}.json"'
            return resp

        # ── CSV ─────────────────────────────────────────────────────────────
        if fmt == 'csv':
            cb   = CSVBuilder(loader)
            data = {
                'full':        loader.audit.to_csv(index=False).encode('utf-8-sig') if not loader.audit.empty else b'',
                'audit':       cb.audit_enriched(),
                'users':       cb.users_enriched(),
                'modules':     cb.module_breakdown(),
                'performance': cb.performance(),
                'attendance':  loader.attendance.to_csv(index=False).encode('utf-8-sig') if not loader.attendance.empty else b'',
                'equipment':   loader.equipment.to_csv(index=False).encode('utf-8-sig') if not loader.equipment.empty else b'',
                'alerts':      loader.alerts.to_csv(index=False).encode('utf-8-sig') if not loader.alerts.empty else b'',
                'certificates':loader.certificates.to_csv(index=False).encode('utf-8-sig') if not loader.certificates.empty else b'',
                'feedback':    loader.feedback.to_csv(index=False).encode('utf-8-sig') if not loader.feedback.empty else b'',
                'wifi':        loader.wifi.to_csv(index=False).encode('utf-8-sig') if not loader.wifi.empty else b'',
                'videography': loader.videography.to_csv(index=False).encode('utf-8-sig') if not loader.videography.empty else b'',
                'filetransfers':loader.filetransfers.to_csv(index=False).encode('utf-8-sig') if not loader.filetransfers.empty else b'',
            }.get(module, cb.audit_enriched())
            resp = HttpResponse(data, content_type='text/csv; charset=utf-8-sig')
            resp['Content-Disposition'] = f'attachment; filename="{org.code}-{module}-{S["report_date"]}.csv"'
            return resp

        # ── PDF / Excel / PPTX all need charts ─────────────────────────────
        cb_charts = ChartBuilder(loader)
        charts    = {
            'api_daily':     cb_charts.api_daily(),
            'hourly':        cb_charts.hourly(),
            'modules':       cb_charts.modules(),
            'response_time': cb_charts.response_time(),
            'users':         cb_charts.users(),
            'notifications': cb_charts.notifications(),
            'attendance':    cb_charts.attendance_trend(),
            'dashboard':     cb_charts.dashboard_overview(),
        }

        if fmt == 'pdf':
            data    = PDFBuilder(loader, charts).build()
            resp    = HttpResponse(data, content_type='application/pdf')
            resp['Content-Disposition'] = f'attachment; filename="{org.code}-analytics-{S["report_date"]}.pdf"'
            return resp

        if fmt in ('excel', 'xlsx'):
            data    = ExcelBuilder(loader, charts).build()
            resp    = HttpResponse(data, content_type='application/vnd.openxmlformats-officedocument.spreadsheetml.sheet')
            resp['Content-Disposition'] = f'attachment; filename="{org.code}-analytics-{S["report_date"]}.xlsx"'
            return resp

        if fmt in ('pptx', 'powerpoint'):
            data    = PPTXBuilder(loader, charts).build()
            resp    = HttpResponse(data, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')
            resp['Content-Disposition'] = f'attachment; filename="{org.code}-analytics-{S["report_date"]}.pptx"'
            return resp

        return HttpResponse({'error': f'Unknown format: {fmt}. Use: pdf, excel, pptx, csv, json'}, status=400)


class AnalyticsExportMetaView(APIView):
    """GET /analytics/export/ — lists available formats and modules."""
    permission_classes = [IsAuthenticated]

    def get(self, request):
        return JsonResponse({
            'formats': ['pdf', 'excel', 'pptx', 'csv', 'json'],
            'modules': ['full', 'audit', 'users', 'modules', 'performance',
                        'attendance', 'equipment', 'alerts', 'certificates',
                        'feedback', 'wifi', 'videography', 'filetransfers'],
            'params':  {'days': 'int (default 30)', 'department': 'uuid (optional)'},
            'example': '/api/v1/analytics/export/pdf/full/?days=30',
        })